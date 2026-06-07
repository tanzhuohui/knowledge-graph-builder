import re
import os
import yaml
import csv


class DocumentChunker:
    IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.webp', '.bmp', '.gif'}

    def __init__(self, config_path=None):
        self.config = self._load_config(config_path)

    def _load_config(self, path):
        if path is None:
            path = os.path.join(
                os.path.dirname(os.path.dirname(__file__)),
                "config", "chunking_config.yaml",
            )
        with open(path, "r", encoding="utf-8") as f:
            return yaml.safe_load(f)["chunking"]

    def chunk_file(self, filepath):
        """根据文件类型分块，支持 md/txt/csv/docx/xlsx/pptx/图片"""
        ext = os.path.splitext(filepath)[1].lower()
        if ext == '.csv':
            return self.chunk_csv(filepath)
        elif ext == '.docx':
            return self.chunk_docx(filepath)
        elif ext == '.xlsx':
            return self.chunk_xlsx(filepath)
        elif ext == '.pptx':
            return self.chunk_pptx(filepath)
        elif ext in self.IMAGE_EXTENSIONS:
            return self.chunk_image(filepath)
        else:
            with open(filepath, "r", encoding="utf-8") as f:
                text = f.read()
            return self.chunk_text(text, os.path.basename(filepath))

    def chunk_csv(self, filepath):
        """将 CSV 文件按行转换为文本块，每行作为一个 chunk"""
        chunks = []
        source = os.path.basename(filepath)
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                reader = csv.DictReader(f)
                headers = reader.fieldnames
                if not headers:
                    return chunks
                
                for i, row in enumerate(reader):
                    # 将每行转换为可读文本描述
                    parts = [f"行 {i+1}:"]
                    for h in headers:
                        val = row.get(h, '').strip()
                        if val:
                            parts.append(f"{h}={val}")
                    text = ", ".join(parts)
                    chunks.append({
                        "id": f"{source}_row_{i}",
                        "source": source,
                        "index": i,
                        "text": text,
                        "word_count": len(text.split()),
                    })
        except Exception as e:
            print(f"  CSV 解析错误 {filepath}: {e}")
        return chunks

    def chunk_docx(self, filepath):
        """将 .docx 文件按段落分块"""
        from docx import Document
        source = os.path.basename(filepath)
        doc = Document(filepath)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        text = "\n\n".join(paragraphs)
        return self.chunk_text(text, source)

    def chunk_xlsx(self, filepath):
        """将 .xlsx 文件每行转为一个文本块"""
        from openpyxl import load_workbook
        source = os.path.basename(filepath)
        wb = load_workbook(filepath, read_only=True, data_only=True)
        chunks = []
        for ws in wb.worksheets:
            sheet_name = ws.title
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            headers = [str(h) if h is not None else "" for h in rows[0]]
            for i, row in enumerate(rows[1:], start=2):
                parts = [f"Sheet={sheet_name}, 行 {i}:"]
                for h, val in zip(headers, row):
                    if val is not None:
                        parts.append(f"{h}={val}")
                text = ", ".join(parts)
                if text.strip():
                    chunks.append({
                        "id": f"{source}_{sheet_name}_row_{i}",
                        "source": source,
                        "index": len(chunks),
                        "text": text,
                        "word_count": len(text.split()),
                    })
        return chunks

    def chunk_pptx(self, filepath):
        """将 .pptx 文件按幻灯片分块，每页为一个 chunk"""
        from pptx import Presentation
        source = os.path.basename(filepath)
        prs = Presentation(filepath)
        chunks = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                text = "\n".join(texts)
                chunks.append({
                    "id": f"{source}_slide_{i}",
                    "source": source,
                    "index": i - 1,
                    "text": text,
                    "word_count": len(text.split()),
                })
        if not chunks:
            chunks.append({
                "id": f"{source}_empty",
                "source": source,
                "index": 0,
                "text": "",
                "word_count": 0,
            })
        return chunks

    def chunk_image(self, filepath):
        """将图片文件转为单 chunk，携带 image_path 供视觉模型处理"""
        source = os.path.basename(filepath)
        return [{
            "id": f"{source}_image_0",
            "source": source,
            "index": 0,
            "text": f"[Image file: {source}]",
            "word_count": 4,
            "image_path": os.path.abspath(filepath),
        }]

    def chunk_text(self, text, source="unknown"):
        strategy = self.config.get("strategy", "fixed")
        if strategy == "fixed":
            return self._fixed_chunk(text, source)
        elif strategy == "sliding_window":
            return self._sliding_window(text, source)
        else:
            return self._fixed_chunk(text, source)

    def _fixed_chunk(self, text, source):
        chunk_size = self.config.get("chunk_size", 1500)
        overlap = self.config.get("overlap", 250)
        separator = self.config.get("separator", "\n\n")

        paragraphs = text.split(separator)
        chunks = []
        current = []
        current_len = 0

        for para in paragraphs:
            para_len = len(para.split())
            if current_len + para_len > chunk_size and current:
                chunks.append(" ".join(current))
                tail = current[-1] if len(current) > 1 else ""
                current = [tail, para] if tail else [para]
                current_len = len(tail.split()) + para_len
            else:
                current.append(para)
                current_len += para_len

        if current:
            chunks.append(" ".join(current))

        return [
            {
                "id": f"{source}_chunk_{i}",
                "source": source,
                "index": i,
                "text": chunk.strip(),
                "word_count": len(chunk.split()),
            }
            for i, chunk in enumerate(chunks)
            if chunk.strip()
        ]

    def _sliding_window(self, text, source):
        chunk_size = self.config.get("chunk_size", 1500)
        overlap = self.config.get("overlap", 250)
        words = text.split()
        chunks = []
        step = chunk_size - overlap

        for i in range(0, len(words), step):
            window = words[i:i + chunk_size]
            if not window:
                break
            chunk_text = " ".join(window)
            chunks.append({
                "id": f"{source}_chunk_{len(chunks)}",
                "source": source,
                "index": len(chunks),
                "text": chunk_text.strip(),
                "word_count": len(window),
            })

        return chunks
